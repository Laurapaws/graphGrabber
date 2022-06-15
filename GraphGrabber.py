# GraphGrabber v1.0 - Laura Moore 2022
# A tool for manipulating specific PDF reports into a specific PowerPoint format
# See https://github.com/Laurapaws/graphGrabber


import glob
import io
import logging
import sys
import os
import time

import platform # Only used to check OS to avoid Win-specific os.startfile
import subprocess
import re # Used in auto import and sort of PDF files
import shutil # Only used for copying some files

import tkinter as tk
from tkinter import filedialog, simpledialog, ttk, Button, Label, Listbox, ANCHOR

import fitz # A PyMuPDF module used here to convert PDF -> PNG
import PIL.Image # For image manipulation
import Pmw # Only used for the Balloon tooltip functionality

from pptx import Presentation #python-pptx is used for manipulating PowerPoints
from pptx.util import Pt

# Set up logging
root_logger= logging.getLogger()
root_logger.setLevel(logging.INFO)
handler = logging.FileHandler('GG.log', 'a', 'utf-8')
handler.setFormatter(logging.Formatter('%(asctime)s | %(name)s | %(message)s'))
root_logger.addHandler(handler)

# All three dictionaries are hardcoded for specific use-cases in EMC Reporting
posDict = {
    # Define coordinates for positioning. Format is test name then test type. e.g. a VT-07 test with a mediumwave plot
    # Tuple Order: Left , Top, Width, Height
    # VT-07
    "VT07MW": (Pt(1), Pt(70), Pt(233), Pt(176)),
    "VT07FM1": (Pt(239), Pt(70), Pt(233), Pt(176)),
    "VT07FM2": (Pt(477), Pt(70), Pt(233), Pt(176)),
    "VT07DAB1AV": (Pt(1), Pt(275), Pt(175), Pt(130)),
    "VT07DAB1RMS": (Pt(180), Pt(275), Pt(175), Pt(130)),
    "VT07DAB2AV": (Pt(360), Pt(275), Pt(175), Pt(130)),
    "VT07DAB2RMS": (Pt(540), Pt(275), Pt(175), Pt(130)),
    # VT-01 3 Metre
    "VT01ThreeVertical": (Pt(1), Pt(85), Pt(358), Pt(270)),
    "VT01ThreeHorizontal": (Pt(360), Pt(85), Pt(358), Pt(270)),
    # VT-12 Single Phase
    "VT12SingleL1": (Pt(1), Pt(85), Pt(358), Pt(270)),
    "VT12SingleN": (Pt(360), Pt(85), Pt(358), Pt(270)),
    # VT-12 Three Phase
    "VT12TripleL1": (Pt(1), Pt(65), Pt(221), Pt(167)),
    "VT12TripleL2": (Pt(222), Pt(65), Pt(221), Pt(167)),
    "VT12TripleL3": (Pt(1), Pt(240), Pt(221), Pt(167)),
    "VT12TripleN": (Pt(222), Pt(240), Pt(221), Pt(167)),
    # VT-15 Electric
    "VT15E16": (Pt(1), Pt(85), Pt(233), Pt(176)),
    "VT15E40": (Pt(239), Pt(85), Pt(233), Pt(176)),
    "VT15E70": (Pt(477), Pt(85), Pt(233), Pt(176)),
    # VT-15 Magnetic Radial and Transverse
    "VT15HR16": (Pt(1), Pt(65), Pt(221), Pt(167)),
    "VT15HR40": (Pt(239), Pt(65), Pt(221), Pt(167)),
    "VT15HR70": (Pt(477), Pt(65), Pt(221), Pt(167)),
    "VT15HT16": (Pt(1), Pt(235), Pt(221), Pt(167)),
    "VT15HT40": (Pt(239), Pt(235), Pt(221), Pt(167)),
    "VT15HT70": (Pt(477), Pt(235), Pt(221), Pt(167)),
}

cropDict = {
    # Define coordinates for cropping. Old refers to the old style PDFs we use that crop in two static positions
    # Left Start, Top Start, Left End, Top End
    "upperOld": ((130, 138, 1000, 800)),
    "lowerOld": ((130, 820, 1000, 1482)),
    "upperOldMagnetic": ((130, 270, 1000, 932)),
}

nameDict = {
    # Takes names of functions and writes their proper name onto the slide
    "VT01Ten": "VT-01 Off-Board Emissions (10m)",
    "VT01Three": "VT-01 Off-Board Emissions (3m)",
    "VT07": "VT-07 On-Board Emissions",
    "VT12Single": "VT-12 Conducted Emissions (Single Phase)",
    "VT12Triple": "VT-12 Conducted Emissions (Three Phase)",
    "VT15Electric": "VT-15 Electric Fields",
    "VT15Magnetic": "VT-15 Magnetic Fields",
}

# Defining global variables
extractedImages = []
croppedImages = []
rejectedList = []
slideCounter = 0
listCounter = 0
cwd = os.getcwd()

currentTime = time.time()
logging.info('************************** Starting GraphGrabber! **************************')


def searchReplace(search_str, repl_str, input, output):
    # Attempts to search and replace on the entire file. Likely needs rewriting to be more robust and not need a template
    # From Stackoverflow
    # The template uses lots of numbers called *1*, *2*... - The function uses these to determine what it should replace
    prs = Presentation(input)
    for slide in prs.slides: # Check in all slides
        for shape in slide.shapes: # Loop objects on slide
            if shape.has_text_frame: # ...Looking for one that contains text
                if (shape.text.find(search_str)) != -1:
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str), str(repl_str)) # Replace what it found with repl_str
                    text_frame.paragraphs[0].runs[0].text = new_text # Honestly, no idea what this line does
    prs.save(output)
    logging.info(search_str + " replaced with " + repl_str)


def extractImages(PDFName, imageFolder):
    # Converts each page of a PDF to an image
    fileName = imageFolder + "/" + PDFName
    doc = fitz.open(fileName)
    zoom = 2  # to increase the resolution
    mat = fitz.Matrix(zoom, zoom)
    noOfPages = doc.pageCount
    for pageNo in range(noOfPages):
        page = doc.load_page(pageNo)  # number of page
        pix = page.get_pixmap(matrix=mat)
        extractedImages.append(pix) # List to store images for current job
        logging.info(
            "Converting " +
            fileName +
            " page " +
            str(pageNo) +
            " to Image")


def cropGraph(targetImg, cropTuple, imName):
    # Crops a specific area of each image according the pixel coordinates defined in the cropTuple
    targetPIL = targetImg.tobytes("PNG")
    im = PIL.Image.open(io.BytesIO(targetPIL)) #BytesIO used to mimic a real file since PIL gets unhappy otherwise
    im1 = im.crop(box=cropTuple) # Crop the newly opened image
    croppedImages.append(im1) # List to store cropped images for current job
    logging.info(imName + " cropped")


def insertImage(oldFileName, newFileName, img, positionTuple, slideNumber):
    # Inserts an image from the croppedImages array into slideNumber using a
    # position from posDict
    prs = Presentation(oldFileName) #oldFileName and newDeckName are distinct to not overwrite the template
    slide = prs.slides[slideNumber] # Currently working on this slide
    left = positionTuple[0] # Position in the slide where this will be inserted
    top = positionTuple[1]
    width = positionTuple[2] # Intended size in inches(?) for the image
    height = positionTuple[3]
    temp = io.BytesIO() # More pretending to be a real file
    img.save(temp, "PNG") # Save the image with a PNG format. JPG quality isn't good enough
    slide.shapes.add_picture(temp, left, top, width, height)
    prs.save(newFileName)
    logging.info(
        "Image inserted with " +
        str(positionTuple) +
        " to " +
        str(slideNumber))


def initialisePowerPoint(emptyDeckName, newDeckName):
    # Sets up the empty, fresh PPTX file

    def resource_path(relative_path):
        try:
            base_path = sys._MEIPASS # Look for the template inside the .exe (lives in temp dir that we can find with this line)
        except Exception:
            base_path = os.path.abspath(".")

        return os.path.join(base_path, relative_path)

    emptyDeckName = resource_path(emptyDeckName)

    emptyDeckName = emptyDeckName + ".pptx" # Hardcoding the extension saves us from having to define it anywhere else (and is lazy)
    newDeckName = newDeckName + ".pptx"
    prs = Presentation(emptyDeckName) # Open the template
    prs.save(newDeckName) # Save it with a new name
    logging.info("Created new PowerPoint: " + newDeckName)

# The first of several hardcoded functions for handling EMC test reports
# Works by extracting all images then cropping them in a specific way depending on the function
def VT07(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName) # Turn each page into an image
    cropGraph(extractedImages[1], cropDict["upperOld"], "MW") # Crop each page in very specific ways and store in a new list
    cropGraph(extractedImages[1], cropDict["lowerOld"], "FM1")
    cropGraph(extractedImages[2], cropDict["upperOld"], "FM2")
    cropGraph(extractedImages[2], cropDict["lowerOld"], "DAB1AV")
    cropGraph(extractedImages[3], cropDict["upperOld"], "DAB1RMS")
    cropGraph(extractedImages[3], cropDict["lowerOld"], "DAB2AV")
    cropGraph(extractedImages[4], cropDict["upperOld"], "DAB2RMS")
    deckName = deckName + ".pptx"
    insertImage(
        deckName,
        deckName,
        croppedImages[0],
        posDict["VT07MW"],
        slideNumber) # Insert all images into their specific locations on the slide
    insertImage(
        deckName,
        deckName,
        croppedImages[1],
        posDict["VT07FM1"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[2],
        posDict["VT07FM2"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[3],
        posDict["VT07DAB1AV"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[4],
        posDict["VT07DAB1RMS"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[5],
        posDict["VT07DAB2AV"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[6],
        posDict["VT07DAB2RMS"],
        slideNumber)
    extractedImages.clear() # Clear both lists
    croppedImages.clear()
    logging.info(">>>>>>>>>>>> Finished VT-07 for " + PDFName)

# The following set of functions are similar with different crops and insertion positions
def VT01Three(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict["upperOld"], "VT01ThreeVertical")
    cropGraph(extractedImages[1], cropDict["lowerOld"], "VT01ThreeHorizontal")
    deckName = deckName + ".pptx"
    insertImage(
        deckName,
        deckName,
        croppedImages[0],
        posDict["VT01ThreeVertical"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[1],
        posDict["VT01ThreeHorizontal"],
        slideNumber,
    )
    extractedImages.clear()
    croppedImages.clear()
    logging.info(">>>>>>>>>>>> Finished VT-01 3m for " + PDFName)


def VT12Single(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict["upperOld"], "VT12SingleL1")
    cropGraph(extractedImages[1], cropDict["lowerOld"], "VT12SingleN")
    deckName = deckName + ".pptx"
    insertImage(
        deckName,
        deckName,
        croppedImages[0],
        posDict["VT12SingleL1"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[1],
        posDict["VT12SingleN"],
        slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    logging.info(">>>>>>>>>>>> Finished VT-12 Single Phase for " + PDFName)


def VT12Triple(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict["upperOld"], "VT12TripleL1")
    cropGraph(extractedImages[1], cropDict["lowerOld"], "VT12TripleL2")
    cropGraph(extractedImages[2], cropDict["upperOld"], "VT12TripleL3")
    cropGraph(extractedImages[2], cropDict["lowerOld"], "VT12TripleN")
    deckName = deckName + ".pptx"
    insertImage(
        deckName,
        deckName,
        croppedImages[0],
        posDict["VT12TripleL1"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[1],
        posDict["VT12TripleL2"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[2],
        posDict["VT12TripleL3"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[3],
        posDict["VT12TripleN"],
        slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    logging.info(">>>>>>>>>>>>  VT-12 Three Phase for " + PDFName)


def VT15Electric(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict["upperOld"], "VT15E16")
    cropGraph(extractedImages[1], cropDict["lowerOld"], "VT15E40")
    cropGraph(extractedImages[2], cropDict["upperOld"], "VT15E70")
    deckName = deckName + ".pptx"
    insertImage(
        deckName,
        deckName,
        croppedImages[0],
        posDict["VT15E16"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[1],
        posDict["VT15E40"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[2],
        posDict["VT15E70"],
        slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    logging.info(">>>>>>>>>>>> Finished VT-15 Electric Field for " + PDFName)


def VT15Magnetic(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict["upperOldMagnetic"], "VT15HR16")
    cropGraph(extractedImages[2], cropDict["upperOld"], "VT15HR40")
    cropGraph(extractedImages[2], cropDict["lowerOld"], "VT15HR70")
    cropGraph(extractedImages[3], cropDict["upperOld"], "VT15HT16")
    cropGraph(extractedImages[3], cropDict["lowerOld"], "VT15HT40")
    cropGraph(extractedImages[4], cropDict["upperOld"], "VT15HT70")
    deckName = deckName + ".pptx"
    insertImage(
        deckName,
        deckName,
        croppedImages[0],
        posDict["VT15HR16"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[1],
        posDict["VT15HR40"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[2],
        posDict["VT15HR70"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[3],
        posDict["VT15HT16"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[4],
        posDict["VT15HT40"],
        slideNumber)
    insertImage(
        deckName,
        deckName,
        croppedImages[5],
        posDict["VT15HT70"],
        slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    logging.info(">>>>>>>>>>>> Finished VT-15 Electric Field for " + PDFName)

# This just sets the slide counter globally
def setSlideCounter(num):
    global slideCounter
    slideCounter = num
    logging.info("Slide counter set to " + str(slideCounter))

# The function responsible for calling other functions (i.e. VT07())
# The parameter reportFunction is used to determine which function it will run
def loopFolder(folderName, deckName, reportFunction):

    directory = folderName # Set dir to current folder we are working on
    global slideCounter # Initialise global variables
    global listCounter
    listCounter = listCounter + 1 #+1 to listCounter so that it skips the header in the GUI listBox
    for file in os.listdir(directory):
        if file.lower().endswith(".pdf"): # Only work on PDFs
            statusMessage = file + " | No Status" # Prepare default status message
            logging.info(
                "Working on slide " +
                str(slideCounter) +
                ", File Name: " +
                file)
            try:
                reportFunction(file, folderName, slideCounter, deckName) # reportFunction() here depends on what you have supplied, e.g. VT07
                searchString = "*" + str(slideCounter) + "*" # Find and replace the title with the file name
                replaceString = (
                    str(file)[:-4] + " | " + (nameDict[str(reportFunction.__name__)])
                )
                searchReplace(
                    searchString,
                    replaceString,
                    deckName + ".pptx",
                    deckName + ".pptx")
                slideCounter = slideCounter + 1
                statusMessage = file + " | Added to slide " + str(slideCounter) # Prepare a status message for the GUI
            except Exception as e:
                statusMessage = file + " | ERROR " + str(e)
                logging.error(e)
            fileList.delete(listCounter) # Remove item from GUI listbox
            fileList.insert(listCounter, statusMessage) # Add in new message with file name and status message
            listCounter = listCounter + 1
            makeProgress() # Update the progress bar
            root.update() # Ensures that the window doesn't feel like it has frozen by updating idle tasks

    logging.info(">>>>>>>>>>>> Finished with folder: " + folderName)


# this is a function to get the selected list box value
def getListboxValue():
    itemSelected = fileList.curselection()
    return itemSelected

# Deprecated button for initialising PowerPoint
def btnInitialisePowerPoint():
    logging.info("Init PP clicked")
    initialisePowerPoint("emptyDeck", "newDeck")

# Create the required folder structure for the hardcoded functions
def btnInitialiseFolders():

    logging.info("Creating Directories")

    def checkCreateDir(dir): # Simple function for checking and making directories
        if os.path.isdir(dir):
            logging.warning(dir + ' already exists')
        else:
            os.mkdir(dir)
            logging.info('CREATED ' + dir)

    checkCreateDir("VT-01 3m") # Run the above func for each folder
    checkCreateDir("VT-07")
    checkCreateDir("VT-12 Single Phase")
    checkCreateDir("VT-12 Three Phase")
    checkCreateDir("VT-15 Electric")
    checkCreateDir("VT-15 Magnetic")
    checkCreateDir("Unsorted PDFs")

# Adds in checks and a prompt for the user to create the required folders instead of doing it silently
def checkFolders():

    dirs = [
        'Unsorted PDFs',
        'VT-01 3m',
        'VT-07',
        'VT-12 Single Phase',
        'VT-12 Three Phase',
        'VT-15 Electric',
        'VT-15 Magnetic']
    missingDirs = [] # Should remain empty if all folders exist

    def checkCreateDir(dir):
        if os.path.isdir(dir):
            logging.warning(dir + ' already exists')
        else:
            os.mkdir(dir)
            logging.info('CREATED ' + dir)

    for dir in dirs:
        if os.path.isdir(dir):
            logging.info(dir + " already exists")
        else:
            missingDirs.append(dir)

    logging.info('Missing Directories: ' + str(missingDirs))

    if len(missingDirs) > 0: # If any of the folders are missing then ask...

        checkFolderMsg = tk.messagebox.askquestion(
            'Missing Directories',
            "Some folders required for GraphGrabber are missing, do you want to create them now?",
            icon='warning')

        if checkFolderMsg == 'yes':
            logging.info("User clicked Yes: Attempting to create directories")
            for dir in missingDirs:
                checkCreateDir(dir)
            tk.messagebox.showinfo(
                'Successfully Created Folders!',
                'Created the following directories:' +
                "\n\n" +
                "\n".join(missingDirs))
        else:
            logging.info("User clicked No: Not creating directories")

# A function that will ask for the file name users want
def askForOutput():
    askPopup = simpledialog.askstring(
        "Output File Name",
        "Please name your output PowerPoint",
        parent=root)
    askPopup = re.sub('[^A-Za-z0-9 ]+', '', askPopup) # Remote all special characters
    return askPopup

# Deletes ALL files in GraphGrabber's folders
def btnClearFolders():

    delList = ['Deleted the following files:']

    def deleteInFolder(dir):
        dir = dir + "/*"
        files = glob.glob(dir)
        for f in files:
            os.remove(f)
            logging.info("DELETED " + str(f)) # Log each deleted file
            delList.append(f) # Add to list of deleted files

    def confirmDel():
        MsgBox = tk.messagebox.askquestion(
            'PDF Deletion',
            "This will delete all files in GraphGrabber's folders:\n VT-01\nVT-07\nVT-12\nVT-15\nUnsorted PDFs\n\nAre you sure?",
            icon='warning')
        if MsgBox == 'yes':
            logging.info("User clicked Yes: Beginning file deletion")
            deleteInFolder("VT-01 3m")
            deleteInFolder("VT-07")
            deleteInFolder("VT-12 Single Phase")
            deleteInFolder("VT-12 Three Phase")
            deleteInFolder("VT-15 Electric")
            deleteInFolder("VT-15 Magnetic")
            deleteInFolder("Unsorted PDFs")
            print(delList)

            tk.messagebox.showinfo(
                'File Deletion Complete',
                "\n".join(delList))

            logging.info("File deletion completed")

            btnCheckFiles()
        else:
            logging.info("User clicked No: Beginning file deletion")

    confirmDel()

# Just opens the current working directory. Differs based on OS
def btnVisitFolders():
    checkFolders()
    path = os.getcwd()
    logging.info("Visiting working directory: " + path)
    if platform.system() == "Windows": # Windows, obviously
        os.startfile(path)
    elif platform.system() == "Darwin": # Mac
        subprocess.Popen(["open", path])
    else: # Other Unix
        subprocess.Popen(["xdg-open", path])

# A function for showing the PDF files in a GUI
def btnCheckFiles():
    logging.info("Checking Files and displaying to user")

    checkFolders()

    global listCounter

    # This function just loops files in a directory and adds them to the tkinter fileList
    def loopInsertList(dir):
        global listCounter

        sectionBreak = "*********** " + dir + " *********** " # A section break in the GUI
        fileList.insert(listCounter, sectionBreak)
        listCounter = listCounter + 1
        localCounter = 0
        for file in os.listdir(dir):
            if file.lower().endswith(".pdf"):
                fileList.insert(listCounter, file)
                localCounter = localCounter + 1 # Local counter counts the files per folder
                listCounter = listCounter + 1 # List counts the entire list including headers

        headerPos = listCounter - localCounter - 1 # Find where the current section header is
        sectionBreak = sectionBreak + "(" + str(localCounter) + " files)" # Add file count onto section break
        fileList.delete(headerPos)
        fileList.insert(headerPos, sectionBreak) # Delete and insert the sectionBreak text into the header index number

    fileList.delete(0, tk.END) # Clear out the GUI file list
    listCounter = 0 # Reset the listCounter
    loopInsertList("VT-01 3m") # Perform the above function for each folder we need to operate on
    loopInsertList("VT-07")
    loopInsertList("VT-12 Single Phase")
    loopInsertList("VT-12 Three Phase")
    loopInsertList("VT-15 Electric")
    loopInsertList("VT-15 Magnetic")
    progessBar["value"] = 0 # Reset the progress bar to zero and set its max value
    progessBar["maximum"] = listCounter - 6

# Main function for running the graph formatter. Primarily uses loopFolder()
def btnGO():
    try:
        checkFolders() # Check that the right folders exist
        logging.info("STARTING JOBS")
        btnCheckFiles() # Display the files to the user (Need a nice clean list)
        setSlideCounter(0) # Reset the slideCounter and listCounter below
        global listCounter
        listCounter = 0
        outputFileName = askForOutput()
        outputFileName = f"{outputFileName} {time.time():.0f}" # Ask user for name for output name. Append unix timestamp
        logging.info('Creating file: ' + outputFileName)
        initialisePowerPoint("emptyDeck", outputFileName) # Prep the initial PowerPoint from template
        loopFolder("VT-01 3m", outputFileName, VT01Three)
        loopFolder("VT-07", outputFileName, VT07)
        loopFolder("VT-12 Single Phase", outputFileName, VT12Single)
        loopFolder("VT-12 Three Phase", outputFileName, VT12Triple)
        loopFolder("VT-15 Electric", outputFileName, VT15Electric)
        loopFolder("VT-15 Magnetic", outputFileName, VT15Magnetic)
        logging.info('>>>>>>>>>>>> JOBS FINISHED')
        tk.messagebox.showinfo(
                'PowerPoint Created',
                'Created ' + outputFileName) # Friendly info box to say its complete
    except Exception as e:
        logging.info('Failed to create deck (no folders?) ' + str(e))
        tk.messagebox.showerror("Error", "Failed to create PowerPoint. Exception has been logged in GG.log") # Error message on failure. Logs to GG.log

# Automatic sorting of files from any folder the user chooses.
def btnAutoSort():
    logging.info('Auto Sort Clicked')

    ceStatus = 1 # Default value for conducted emissions tests is set to single phase

    try:
        if getphaseListValue()[0] == 1: # Check if the user has selected single or three phase to sort
            ceStatus = 3
            logging.info(str(ceStatus) + ' phase')
        else:
            ceStatus = 1
            logging.info(str(ceStatus) + ' phase')
    except BaseException:
        logging.info('Defaulting to single phase')
        # Needs rewriting as a message to ask user (only if VT12 exists)

    def regexCopy(file, dir, destination): # Copies the file specified later
        fileToCopy = dir + '/' + file
        shutil.copy(fileToCopy, destination)
        logging.info('COPIED to ' + destination + ': ' + fileToCopy)

    try:
        dir = filedialog.askdirectory()  # Returns opened path as str

        for file in os.listdir(dir):
            if file.lower().endswith(".pdf"):
                if re.search('REESS', file, flags=re.I):# Attempt to sort PDFs into folders based on names
                    regexCopy(file, dir, 'VT-01 3m')
                elif re.search('NB', file, flags=re.I):
                    regexCopy(file, dir, 'VT-01 3m')
                elif re.search('BB', file, flags=re.I):
                    regexCopy(file, dir, 'VT-01 3m')
                elif re.search('e.field', file, flags=re.I):
                    regexCopy(file, dir, 'VT-15 Electric')
                elif re.search('H.Field', file, flags=re.I):
                    regexCopy(file, dir, 'VT-15 Magnetic')
                elif re.search('(?<!i)(?<!n)CE', file, flags=re.I):
                    if ceStatus == 1:
                        regexCopy(file, dir, 'VT-12 Single Phase')
                    else:
                        regexCopy(file, dir, 'VT-12 Three Phase')
                else:
                    regexCopy(file, dir, 'Unsorted PDFs')
    except Exception as e:
        logging.info('Auto sort failed with ' + str(e))

    btnCheckFiles()


# This is a function which increases the progress bar value by the given
# increment amount
def makeProgress():
    progessBar["value"] = progessBar["value"] + 1
    root.update_idletasks()


# this is a function to get the fileList list box value
# fileList is the main GUI box that shows the PDFs being worked on
def getfileListValue():
    itemSelected = fileList.curselection()
    return itemSelected

# this is a function to get the phaseList list box value
# phaseList is the box for conducted emissions autoselect sorting
def getphaseListValue():
    itemSelected = phaseList.curselection()
    return itemSelected


root = tk.Tk()
# Create the main window
root.geometry("550x850")
root.configure(background="#34495e")
root.title("GraphGrabber v1")
root.iconbitmap('ICOLogo.ico')

# Initialise Python Megawidgets
Pmw.initialise(root)

# Clear Folders Button
wgtClearFolders = Button(
    root,
    text="Clear\nFolders",
    fg="#f1c40f",
    bg='#c0392b',
    font=("Helvetica", 13, "normal"),
    command=btnClearFolders,
)
wgtClearFolders.place(x=475, y=10) # Split into multiple parts to allow it to be referenced later for the tooltip

tipName = Pmw.Balloon(root)
tipName.bind(
    wgtClearFolders,
    '''This will delete everything in the folders created by GraphGrabber
Do not have anything stored in here that you want to keep!''')

# Current Working Directory Label
Label(
    root,
    text=cwd,
    fg="#f1c40f",
    bg="#34495e",
    wraplength=530,
    justify="left",
    font=("courier", 10, "normal"),
).place(x=10, y=64)


# Go to Current Working Directory Button
wgtVisitFolders = Button(
    root,
    text="Open Working\nDirectory",
    bg="#bdc3c7",
    fg="#2c3e50",
    font=("Helvetica", 13, "normal"),
    command=btnVisitFolders,
)
wgtVisitFolders.place(x=10, y=10)

tipName = Pmw.Balloon(root)
tipName.bind(
    wgtVisitFolders,
    '''This will open the current working directory as displayed above.
By default this is the folder where GraphGrabber.exe lives
Move the .exe somewhere else to change this folder.''')

# Check Files Button
wgtCheckFiles = Button(
    root,
    text="Check\nFiles",
    bg="#3498db",
    fg="#ecf0f1",
    font=("Helvetica", 13, "normal"),
    command=btnCheckFiles,
)
wgtCheckFiles.place(x=328, y=10)

tipName = Pmw.Balloon(root)
tipName.bind(
    wgtCheckFiles,
    '''This will scan through the current directory
Use it to check that you have all files in the right places
Needs the folder structure created with the Initialise Folders button''')

# Auto Sort Button
wgtAutoSort = Button(
    root,
    text="Auto\nSort",
    bg="#bdc3c7",
    fg="#2c3e50",
    font=("Helvetica", 13, "normal"),
    command=btnAutoSort,
)
wgtAutoSort.place(x=140, y=10)

tipAutoSort = Pmw.Balloon(root)
tipAutoSort.bind(wgtAutoSort, '''Select the folder containing report PDFs
GraphGrabber will attempt to sort them into its folders for you
Unsorted files will go into the Unsorted PDFs folder
Select single or three phase below for conducted emissions''')

# Create Deck Button
wgtGO = Button(
    root,
    text="Create\nDeck!",
    bg="#1abc9c",
    fg="#ecf0f1",
    font=("Helvetica", 13, "normal"),
    command=btnGO,
)
wgtGO.place(x=401, y=10)

tipGO = Pmw.Balloon(root)
tipGO.bind(wgtGO, '''Starts creating the Powerpoint
FIle name will be your output with the Unix epoch
Press Initialise Folders to create the right folder structure
Press Check Files so you know what the program will operate on''')

# Progress Bar
progessBar_style = ttk.Style()
progessBar_style.theme_use("clam")
progessBar_style.configure(
    "progessBar.Horizontal.TProgressbar",
    foreground="#00CD00",
    background="#00CD00")
progessBar = ttk.Progressbar(
    root,
    style="progessBar.Horizontal.TProgressbar",
    orient="horizontal",
    length=530,
    mode="determinate",
    maximum=100,
    value=0,
)
progessBar.place(x=10, y=105)


# File List
fileList = Listbox(
    root, bg="#bdc3c7", font=("Helvetica", 10, "normal"), width=75, height=41
)
fileList.place(x=10, y=130)

# Double click on the listbox to copy that entry
def listbox_copy(event):
    root.clipboard_clear()
    selected = fileList.get(ANCHOR)
    root.clipboard_append(selected)

fileList.bind('<Double-Button-1>', listbox_copy)

# VT-12 Phase List
phaseList = Listbox(
    root,  bg="#bdc3c7", font=("Helvetica", 11, "normal"), width=15, height=2
)
phaseList.insert('0', 'Single-Phase CE')
phaseList.insert('1', 'Three-Phase CE')
phaseList.place(x=195, y=15)

# Run folder check on startup
checkFolders()

# Run tkinter window
root.mainloop()


logging.info("************************** Window Closing... **************************")
